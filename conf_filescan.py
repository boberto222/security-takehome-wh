import os
import shutil
from docx import Document
from pptx import Presentation
import pandas as pd
import csv
import datetime

#For dependencies - use pip3 to install python-docx, csv, datetime, pandas, and python-pptx if executing in bash or zsh terminal prior

keywords = ["Confidential", "Analysis", "Internal", "Private", "Restricted"]
filestore = "/Users/williamherrera/Downloads/DataLakeFiles"
conf_dir = "/Users/williamherrera/Downloads/DataLakeFiles/Confidential_Files"
pub_dir = "/Users/williamherrera/Downloads/DataLakeFiles/Public_Files"
report_dir = "/Users/williamherrera/Downloads/DataLakeFiles/Scan_Reports"

#Update directories to your own configuration. keywords matrix is used for determining if confidential or not"

def check_keywords(text):
    for a in keywords:
        if a in text:
            return True
    return False

#Definition for keyword review in the text of each document

def move_file(file, is_confidential):
    if is_confidential:
        shutil.move(os.path.join(filestore, file), os.path.join(conf_dir, os.path.basename(file)))
        csv_writer.writerow([file, 'Confidential', datetime.datetime.now()])
    else:
        shutil.move(os.path.join(filestore, file), os.path.join(pub_dir, os.path.basename(file)))
        csv_writer.writerow([file, 'Public', datetime.datetime.now()])

#Definition on if the text Definition is true, move to either conf_dir or pub_dir. When it is moved, an entry is added to the CSV that is created.

csv_file = open(str(report_dir) + "/Report_" + str(datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")) + '.csv', 'w', newline='')
csv_writer = csv.writer(csv_file)
csv_writer.writerow(['Filename', 'Confidential or Public', 'Time Completed'])

#Upon execution, create a CSV with the three column headers, named "Report" with the systemtime of the execution in title.


for file in os.listdir(filestore):
    if file.endswith(".docx"):
        conf_doc = Document(os.path.join(filestore, file))
        text = "\n".join([paragraph.text for paragraph in conf_doc.paragraphs])
        is_confidential = check_keywords(text)
        move_file(file, is_confidential)
    
    elif file.endswith(".pptx"):
        conf_pres = Presentation(os.path.join(filestore, file))
        text = ""
        for slide in conf_pres.slides:
        	for shape in slide.shapes:
        		if hasattr(shape, "text"):
        			text += shape.text
        is_confidential = check_keywords(text)
        move_file(file, is_confidential)
    
    elif file.endswith(".xlsx"):
        conf_book = os.path.join(filestore, file)
        df = pd.read_excel(conf_book)
        text = df.to_string()
        is_confidential = check_keywords(text)
        move_file(file, is_confidential)
